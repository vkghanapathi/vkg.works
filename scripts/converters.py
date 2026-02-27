"""
converters.py — Format-specific article converters for vkg.works
Supports: .md, .html, .docx, .pptx, .pdf, .txt
All I/O is UTF-8 with ensure_ascii=False where applicable.
VijayaDV PUA Unicode characters are never normalised or stripped.
"""
from __future__ import annotations
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional
from datetime import date as date_type

import frontmatter
from markdown_it import MarkdownIt
from dateutil.parser import parse as parse_date


@dataclass
class ArticleData:
    slug: str
    section: str
    title: str
    date: Optional[str]          # ISO format "YYYY-MM-DD"
    date_display: Optional[str]  # Human-readable
    excerpt: str
    author: str
    body_html: str
    is_pdf: bool = False
    pdf_filename: Optional[str] = None
    audio_file: Optional[str] = None
    youtube_url: Optional[str] = None
    notation_pdf: Optional[str] = None
    video_file: Optional[str] = None
    coverage_type: Optional[str] = None   # "press" or "event"
    source_url: Optional[str] = None
    photos: list = field(default_factory=list)
    pdf_file: Optional[str] = None        # For books
    status: Optional[str] = None          # For projects
    category: Optional[str] = None        # Free-text type (administrative/technical/etc.)
    ref: Optional[str] = None             # Permanent section ref e.g. A·0001
    featured: bool = False                # Pin to home page featured column
    extracted_images: list = field(default_factory=list)  # [(filename, bytes), …] from DOCX


_md = MarkdownIt()
_DATE_PREFIX = re.compile(r'^(\d{4}-\d{2}-\d{2})-(.+)$')


def _parse_filename(stem: str) -> tuple[Optional[str], str]:
    """Extract (date_str, slug) from filename stem."""
    m = _DATE_PREFIX.match(stem)
    if m:
        return m.group(1), m.group(2)
    return None, stem


def _fmt_date(date_str: Optional[str]) -> Optional[str]:
    if not date_str:
        return None
    try:
        d = parse_date(date_str)
        return d.strftime("%-d %B %Y") if hasattr(d, 'strftime') else date_str
    except Exception:
        return date_str


def _truncate(text: str, length: int = 200) -> str:
    text = text.strip()
    if len(text) <= length:
        return text
    return text[:length].rsplit(' ', 1)[0] + '…'


def _first_paragraph(text: str, max_chars: int = 220) -> str:
    """Extract first 1-2 sentences from plain text."""
    # Strip markdown headings and blank lines
    lines = [l.strip() for l in text.splitlines()
             if l.strip() and not l.strip().startswith('#')]
    plain = ' '.join(lines)
    # Split on sentence-ending punctuation followed by space
    sentences = re.split(r'(?<=[.!?।])\s+', plain)
    result = ''
    for sent in sentences:
        if len(result) + len(sent) <= max_chars:
            result = (result + ' ' + sent).strip()
        else:
            break
        if len(result) >= 80:   # at least one decent sentence
            break
    return result or _truncate(plain, max_chars)


def _youtube_to_embed(url: str) -> Optional[str]:
    """Convert a YouTube watch/youtu.be URL to an embed URL."""
    if not url:
        return None
    # Already embed URL
    if 'youtube.com/embed' in url:
        return url
    # youtu.be/ID
    m = re.match(r'https?://youtu\.be/([^?&]+)', url)
    if m:
        return f"https://www.youtube.com/embed/{m.group(1)}"
    # youtube.com/watch?v=ID
    m = re.match(r'https?://(?:www\.)?youtube\.com/watch\?.*v=([^&]+)', url)
    if m:
        return f"https://www.youtube.com/embed/{m.group(1)}"
    # live stream
    m = re.match(r'https?://(?:www\.)?youtube\.com/live/([^?&]+)', url)
    if m:
        return f"https://www.youtube.com/embed/{m.group(1)}"
    return url


# ─────────────────────────────────────────────
# Format converters
# ─────────────────────────────────────────────

def convert_markdown(path: Path, section: str) -> ArticleData:
    date_str, slug = _parse_filename(path.stem)
    post = frontmatter.load(str(path))
    meta = post.metadata

    title = str(meta.get('title', slug.replace('-', ' ').title()))
    if not date_str:
        date_str = str(meta.get('date', '')) or None
    excerpt = str(meta.get('excerpt', '')) or _first_paragraph(post.content)
    author = str(meta.get('author', 'Dr. Vamshi Krishna Ghanapāṭhī'))
    body_html = _md.render(post.content)

    youtube_raw = meta.get('youtube_url', '')
    return ArticleData(
        slug=slug,
        section=section,
        title=title,
        date=date_str,
        date_display=_fmt_date(date_str),
        excerpt=_truncate(excerpt),
        author=author,
        body_html=body_html,
        audio_file=meta.get('audio_file'),
        youtube_url=_youtube_to_embed(youtube_raw) if youtube_raw else None,
        notation_pdf=meta.get('notation_pdf'),
        video_file=meta.get('video_file'),
        coverage_type=meta.get('type'),
        source_url=meta.get('source_url'),
        pdf_file=meta.get('pdf_file'),
        status=meta.get('status'),
        category=meta.get('category') or None,
        featured=bool(meta.get('featured', False)),
    )


def convert_html(path: Path, section: str) -> ArticleData:
    date_str, slug = _parse_filename(path.stem)
    content = path.read_text(encoding='utf-8')
    # Extract <title> if present
    m = re.search(r'<title>(.*?)</title>', content, re.IGNORECASE | re.DOTALL)
    title = m.group(1).strip() if m else slug.replace('-', ' ').title()
    # Strip <html><head> etc if present — just keep body content
    body_m = re.search(r'<body[^>]*>(.*?)</body>', content, re.IGNORECASE | re.DOTALL)
    body_html = body_m.group(1) if body_m else content
    excerpt = _first_paragraph(re.sub(r'<[^>]+>', '', body_html))
    return ArticleData(
        slug=slug, section=section, title=title,
        date=date_str, date_display=_fmt_date(date_str),
        excerpt=_truncate(excerpt), author='Dr. Vamshi Krishna Ghanapāṭhī',
        body_html=body_html,
    )


def convert_docx(path: Path, section: str) -> ArticleData:
    from docx import Document
    from docx.oxml.ns import qn as _qn
    date_str, slug = _parse_filename(path.stem)
    doc = Document(str(path))

    title = slug.replace('-', ' ').title()
    paragraphs_html = []
    first_text_para = None
    extracted_images: list = []
    img_idx = 0

    for para in doc.paragraphs:
        # Extract any inline images from this paragraph
        for drawing in para._p.findall('.//' + _qn('w:drawing')):
            for blip in drawing.findall('.//' + _qn('a:blip')):
                r_embed = blip.get(_qn('r:embed'))
                if r_embed:
                    try:
                        img_part = doc.part.related_parts[r_embed]
                        ext = img_part.partname.rpartition('.')[-1].lower() or 'png'
                        if ext not in ('png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp'):
                            ext = 'png'
                        fname = f'{slug}-img-{img_idx:03d}.{ext}'
                        extracted_images.append((fname, img_part.blob))
                        paragraphs_html.append(
                            f'<figure class="article-figure">'
                            f'<img src="{fname}" alt="Figure {img_idx + 1}" loading="lazy">'
                            f'</figure>'
                        )
                        img_idx += 1
                    except Exception:
                        pass

        text = para.text.strip()
        if not text:
            continue
        style = para.style.name if para.style else ''
        if style in ('Title', 'Heading 1') and title == slug.replace('-', ' ').title():
            title = text
        elif style.startswith('Heading'):
            level = '2'
            m = re.search(r'\d', style)
            if m:
                level = min(int(m.group()), 4)
            paragraphs_html.append(f'<h{level}>{text}</h{level}>')
        else:
            # Build inline HTML for runs (bold, italic)
            runs_html = ''
            for run in para.runs:
                t = run.text
                if not t:
                    continue
                if run.bold and run.italic:
                    t = f'<strong><em>{t}</em></strong>'
                elif run.bold:
                    t = f'<strong>{t}</strong>'
                elif run.italic:
                    t = f'<em>{t}</em>'
                runs_html += t
            if runs_html:
                paragraphs_html.append(f'<p>{runs_html}</p>')
                if first_text_para is None:
                    first_text_para = text

    body_html = '\n'.join(paragraphs_html)
    excerpt = _truncate(first_text_para or '')

    # Check for a .md sidecar that provides frontmatter metadata
    sidecar = path.with_suffix('.md')
    status = None
    category = None
    featured = False
    if sidecar.exists():
        try:
            post = frontmatter.load(str(sidecar))
            meta = post.metadata
            if meta.get('title'):
                title = str(meta['title'])
            if meta.get('date'):
                date_str = str(meta['date'])
            if meta.get('author'):
                pass  # use sidecar author below
            status = meta.get('status') or None
            category = meta.get('category') or None
            featured = bool(meta.get('featured', False))
            if meta.get('excerpt'):
                excerpt = str(meta['excerpt'])
        except Exception:
            pass

    return ArticleData(
        slug=slug, section=section, title=title,
        date=date_str, date_display=_fmt_date(date_str),
        excerpt=excerpt, author='Dr. Vamshi Krishna Ghanapāṭhī',
        body_html=body_html,
        status=status, category=category, featured=featured,
        extracted_images=extracted_images,
    )


def convert_pptx(path: Path, section: str) -> ArticleData:
    from pptx import Presentation
    date_str, slug = _parse_filename(path.stem)
    prs = Presentation(str(path))

    title = slug.replace('-', ' ').title()
    sections_html = []
    first_body_text = None

    for i, slide in enumerate(prs.slides):
        slide_title = ''
        slide_body = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.shape_type == 13:  # Picture
                continue
            is_title = (hasattr(shape, 'placeholder_format') and
                        shape.placeholder_format is not None and
                        shape.placeholder_format.type in (1, 15))  # CENTER_TITLE or TITLE
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if is_title:
                slide_title = text
                if i == 0:
                    title = text
            else:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_body.append(f'<p>{t}</p>')
                        if first_body_text is None and not is_title:
                            first_body_text = t

        slide_html = f'<section class="slide">'
        if slide_title:
            slide_html += f'<h2>{slide_title}</h2>'
        slide_html += '\n'.join(slide_body) + '</section>'
        sections_html.append(slide_html)

    body_html = '\n'.join(sections_html)
    excerpt = _truncate(first_body_text or '')
    return ArticleData(
        slug=slug, section=section, title=title,
        date=date_str, date_display=_fmt_date(date_str),
        excerpt=excerpt, author='Dr. Vamshi Krishna Ghanapāṭhī',
        body_html=body_html,
    )


def convert_pdf(path: Path, section: str) -> ArticleData:
    import pdfplumber
    date_str, slug = _parse_filename(path.stem)
    title = slug.replace('-', ' ').title()
    excerpt = ''
    try:
        with pdfplumber.open(str(path)) as pdf:
            if pdf.pages:
                text = pdf.pages[0].extract_text() or ''
                lines = [l.strip() for l in text.splitlines() if l.strip()]
                if lines:
                    title = lines[0] if len(lines[0]) < 120 else title
                excerpt = _first_paragraph('\n'.join(lines[1:]) if len(lines) > 1 else text)
    except Exception:
        pass

    return ArticleData(
        slug=slug, section=section, title=title,
        date=date_str, date_display=_fmt_date(date_str),
        excerpt=_truncate(excerpt), author='Dr. Vamshi Krishna Ghanapāṭhī',
        body_html='', is_pdf=True, pdf_filename=path.name,
    )


def convert_txt(path: Path, section: str) -> ArticleData:
    date_str, slug = _parse_filename(path.stem)
    content = path.read_text(encoding='utf-8')
    lines = [l.rstrip() for l in content.splitlines()]
    non_empty = [l for l in lines if l.strip()]

    title = slug.replace('-', ' ').title()
    if non_empty and len(non_empty[0]) < 100 and not non_empty[0].endswith('.'):
        title = non_empty[0]
        body_lines = lines[1:]
    else:
        body_lines = lines

    body_text = '\n'.join(body_lines)
    excerpt = _first_paragraph('\n'.join(non_empty[1:]) if len(non_empty) > 1 else content)
    body_html = f'<pre style="white-space:pre-wrap; font-family:inherit;">{body_text}</pre>'

    return ArticleData(
        slug=slug, section=section, title=title,
        date=date_str, date_display=_fmt_date(date_str),
        excerpt=_truncate(excerpt), author='Dr. Vamshi Krishna Ghanapāṭhī',
        body_html=body_html,
    )


def convert_audio(path: Path, section: str, md_path: Optional[Path] = None) -> ArticleData:
    """For standalone .mp3/.m4a files."""
    date_str, slug = _parse_filename(path.stem)
    title = slug.replace('-', ' ').title()
    excerpt = ''
    author = 'Dr. Vamshi Krishna Ghanapāṭhī'
    body_html = ''

    if md_path and md_path.exists():
        post = frontmatter.load(str(md_path))
        meta = post.metadata
        title = str(meta.get('title', title))
        excerpt = str(meta.get('excerpt', ''))
        author = str(meta.get('author', author))
        body_html = _md.render(post.content) if post.content.strip() else ''
        if not date_str:
            date_str = str(meta.get('date', '')) or None

    suffix = path.suffix.lower()
    mime = 'audio/mpeg' if suffix == '.mp3' else 'audio/mp4'

    return ArticleData(
        slug=slug, section=section, title=title,
        date=date_str, date_display=_fmt_date(date_str),
        excerpt=_truncate(excerpt), author=author, body_html=body_html,
        audio_file=path.name,
        # Store mime type in notation_pdf field (repurposed as mime hint)
        notation_pdf=mime,
    )


def convert_video_md(path: Path, section: str) -> ArticleData:
    """Video entries are always Markdown files with youtube_url or video_file frontmatter."""
    data = convert_markdown(path, section)
    if data.youtube_url:
        data.youtube_url = _youtube_to_embed(data.youtube_url)
    return data


def convert_coverage(path: Path, section: str) -> ArticleData:
    """Coverage items: either a .md file or a folder with index.md + images."""
    if path.is_dir():
        md_file = path / 'index.md'
        data = convert_markdown(md_file, section) if md_file.exists() else ArticleData(
            slug=path.name, section=section,
            title=path.name.replace('-', ' ').title(),
            date=None, date_display=None, excerpt='', author='',
            body_html='',
        )
        # Collect images in the folder
        image_exts = {'.jpg', '.jpeg', '.png', '.gif', '.webp'}
        photos = [
            {'url': f'/coverage/{path.name}/{f.name}', 'name': f.stem}
            for f in sorted(path.iterdir())
            if f.suffix.lower() in image_exts
        ]
        data.photos = photos
        # Use folder slug (strip date prefix)
        date_str, slug = _parse_filename(path.name)
        data.slug = slug
        if not data.date:
            data.date = date_str
            data.date_display = _fmt_date(date_str)
        return data
    else:
        return convert_markdown(path, section)


CONVERTERS = {
    '.md':   convert_markdown,
    '.html': convert_html,
    '.htm':  convert_html,
    '.docx': convert_docx,
    '.pptx': convert_pptx,
    '.pdf':  convert_pdf,
    '.txt':  convert_txt,
}
